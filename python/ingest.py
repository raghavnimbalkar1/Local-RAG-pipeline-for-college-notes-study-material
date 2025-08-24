#!/usr/bin/env python3
"""
Ingest PDFs and PPTX files from data/raw/** and write token-based chunks to data/processed/chunks.jsonl

Each JSONL line has fields: id, text, source, page
- id format: "<subject>/<filename>#p<page>#c<chunk>"
- source is the relative path under data/raw
- page is 1-based page (PDF) or slide number (PPTX)

Usage:
  python python/ingest.py \
    --raw-dir data/raw \
    --out data/processed/chunks.jsonl \
    --model sentence-transformers/all-MiniLM-L6-v2 \
    --max-tokens 700 --overlap 100
"""

from __future__ import annotations
import argparse
import json
from pathlib import Path
from typing import Iterable, List, Tuple

from pypdf import PdfReader
from pptx import Presentation
from transformers import AutoTokenizer, AutoConfig

# OCR fallback
import pytesseract
from pdf2image import convert_from_path

SUPPORTED_EXTS = {".pdf", ".pptx"}


def iter_files(root: Path) -> Iterable[Path]:
    for p in root.rglob("*"):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS:
            yield p


def ocr_pdf_page(path: Path, page_num: int) -> str:
    """Run OCR on a single PDF page (1-based index)."""
    try:
        images = convert_from_path(str(path), first_page=page_num, last_page=page_num)
        if not images:
            return ""
        return pytesseract.image_to_string(images[0])
    except Exception as e:
        print(f"[WARN] OCR failed on {path} page {page_num}: {e}")
        return ""


def read_pdf(path: Path) -> List[Tuple[int, str]]:
    pages: List[Tuple[int, str]] = []
    try:
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages, start=1):
            try:
                txt = page.extract_text() or ""
            except Exception:
                txt = ""
            txt = txt.replace("\r", "\n")

            if not txt.strip():
                txt = ocr_pdf_page(path, i)

            pages.append((i, txt))
    except Exception as e:
        print(f"[WARN] Failed to read PDF: {path} ({e})")
    return pages


def _shape_text(shape) -> str:
    out = []
    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            out.append(shape.text)
        if getattr(shape, "has_table", False) and shape.table:
            for r in shape.table.rows:
                for c in r.cells:
                    out.append(c.text)
    except Exception:
        pass
    return "\n".join([s for s in out if s])


def read_pptx(path: Path) -> List[Tuple[int, str]]:
    slides: List[Tuple[int, str]] = []
    try:
        pres = Presentation(str(path))
        for i, slide in enumerate(pres.slides, start=1):
            texts = []
            for shp in slide.shapes:
                t = _shape_text(shp)
                if t:
                    texts.append(t)
            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    texts.append(slide.notes_slide.notes_text_frame.text)
            except Exception:
                pass
            body = "\n".join([t for t in texts if t])
            slides.append((i, body))
    except Exception as e:
        print(f"[WARN] Failed to read PPTX: {path} ({e})")
    return slides


def chunk_tokens(text: str, tokenizer, max_tokens: int, overlap: int) -> List[str]:
    if not text.strip():
        return []
    ids = tokenizer.encode(text, add_special_tokens=False, truncation=True, max_length=tokenizer.model_max_length)
    if not ids:
        return []
    chunks: List[str] = []
    stride = max(1, max_tokens - overlap)
    for start in range(0, len(ids), stride):
        end = min(start + max_tokens, len(ids))
        window = ids[start:end]
        if not window:
            break
        chunk_text = tokenizer.decode(window, skip_special_tokens=True)
        if chunk_text.strip():
            chunks.append(chunk_text)
        if end >= len(ids):
            break
    return chunks


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--raw-dir", type=Path, default=Path("data/raw"))
    ap.add_argument("--out", type=Path, default=Path("data/processed/chunks.jsonl"))
    ap.add_argument("--model", type=str, default="sentence-transformers/all-MiniLM-L6-v2")
    ap.add_argument("--max-tokens", type=int, default=700)
    ap.add_argument("--overlap", type=int, default=100)
    args = ap.parse_args()

    raw_dir: Path = args.raw_dir
    out_path: Path = args.out
    out_path.parent.mkdir(parents=True, exist_ok=True)

    tokenizer = AutoTokenizer.from_pretrained(args.model)
    config = AutoConfig.from_pretrained(args.model)

    # Respect model maximum sequence length
    model_max_len = getattr(config, "max_position_embeddings", 512)
    if args.max_tokens > model_max_len:
        print(f"[INFO] Reducing max_tokens from {args.max_tokens} → {model_max_len} (model limit)")
        args.max_tokens = model_max_len

    total_files = 0
    total_pages = 0
    total_chunks = 0

    with out_path.open("w", encoding="utf-8") as fout:
        for fpath in iter_files(raw_dir):
            rel = fpath.relative_to(raw_dir).as_posix()
            ext = fpath.suffix.lower()
            total_files += 1

            if ext == ".pdf":
                pages = read_pdf(fpath)
            elif ext == ".pptx":
                pages = read_pptx(fpath)
            else:
                continue

            if all((t or "").strip() == "" for _, t in pages):
                print(f"[WARN] No extractable text (even OCR failed): {rel}")

            for page_num, page_text in pages:
                total_pages += 1
                chunks = chunk_tokens(page_text, tokenizer, args.max_tokens, args.overlap)
                if not chunks and page_text.strip():
                    chunks = [page_text.strip()]
                for ci, chunk in enumerate(chunks, start=1):
                    cid = f"{rel}#p{page_num}#c{ci}"
                    rec = {
                        "id": cid,
                        "text": chunk,
                        "source": rel,
                        "page": page_num,
                    }
                    fout.write(json.dumps(rec, ensure_ascii=False) + "\n")
                    total_chunks += 1

    print(
        f"✅ Done. Files: {total_files} | Pages/Slides: {total_pages} | Chunks: {total_chunks}\n"
        f"Output → {out_path}"
    )


if __name__ == "__main__":
    main()
